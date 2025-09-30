import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def extract_textbox_properties(pptx_path: str, output_json: str = None):
    """
    Extract text box properties (position, size, font, color, text, alignment, margins) from a PPTX file
    and return/save them as JSON.
    """

    prs = Presentation(pptx_path)

    fields = {}
    field_index = 1

    # Helper: convert EMU → px
    def emu_to_px(emu):
        return int(emu / 914400 * 96)

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            # Position and size
            x, y = emu_to_px(shape.left), emu_to_px(shape.top)

            # Margins (padding)
            tf = shape.text_frame
            margins = {
                "left": emu_to_px(tf.margin_left),
                "right": emu_to_px(tf.margin_right),
                "top": emu_to_px(tf.margin_top),
                "bottom": emu_to_px(tf.margin_bottom),
            }

            # Alignment – take from first paragraph
            alignment = "left"
            if tf.paragraphs and tf.paragraphs[0].alignment is not None:
                align_map = {
                    PP_ALIGN.LEFT: "left",
                    PP_ALIGN.CENTER: "center",
                    PP_ALIGN.RIGHT: "right",
                }
                alignment = align_map.get(tf.paragraphs[0].alignment, "left")

            # Font defaults
            font_name = None
            font_size = None
            font_color = "black"
            bold = False
            italic = False

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    if run.font:
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            font_size = int(run.font.size.pt)
                        if run.font.bold:
                            bold = True
                        if run.font.italic:
                            italic = True
                        try:
                            rgb = run.font.color.rgb
                            if rgb:
                                font_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                        except AttributeError:
                            pass

            fields[f"field{field_index}"] = {
                "pos": [x, y],
                "margins": margins,
                "font": font_name or "Arial",
                "font_size": font_size or 24,
                "bold": bold,
                "italic": italic,
                "fill": font_color,
                "align": alignment,
                "sample_text": text,
            }

            field_index += 1

    if output_json:
        with open(output_json, "w", encoding="utf-8") as f:
            json.dump(fields, f, indent=2, ensure_ascii=False)

    return fields


if __name__ == "__main__":
    pptx_file = "template-filled.pptx"  # Replace with your actual file
    output = "fields_config.json"

    config = extract_textbox_properties(pptx_file, output_json=output)
    print(json.dumps(config, indent=2, ensure_ascii=False))